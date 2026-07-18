"""Shared style tokens and shape helpers for Verbatim's presentation decks.

Design language (colors, type, icon geometry) lives here once so every deck
script (`build_deck.py`, `build_deck_week2.py`, ...) draws from the same
visual system instead of duplicating ~400 lines of python-pptx boilerplate.
"""

from __future__ import annotations

from pathlib import Path

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.presentation import Presentation
from pptx.slide import Slide
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).parent
ICON_PATH = HERE.parent / "addon" / "icon.png"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BLUE = RGBColor(0x42, 0x85, 0xF4)
INK = RGBColor(0x20, 0x21, 0x24)
SECONDARY = RGBColor(0x5F, 0x63, 0x68)
BG = RGBColor(0xF8, 0xF9, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DIVIDER = RGBColor(0xDA, 0xDC, 0xE0)

CATEGORY_COLORS = {
    "tone_drift": RGBColor(0x8E, 0x24, 0xAA),
    "information_hierarchy": RGBColor(0x1E, 0x88, 0xE5),
    "cta_cadence": RGBColor(0x00, 0x89, 0x7B),
    "readability": RGBColor(0xF4, 0x51, 0x1E),
    "formatting_and_style": RGBColor(0x6D, 0x4C, 0x41),
    "channel_constraints": RGBColor(0x39, 0x49, 0xAB),
    "banned_words_and_competitors": RGBColor(0xC6, 0x28, 0x28),
}

HEAD_FONT = "Google Sans"
BODY_FONT = "Roboto"


def set_run(run, *, font=BODY_FONT, size=18, color=INK, bold=False, italic=False):
    """Apply consistent font styling to a text run."""
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def add_background(slide: Slide, color=WHITE):
    """Fill the slide background."""
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_W, SLIDE_H)
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    rect.shadow.inherit = False
    # Push the background behind everything else added to the slide.
    sp_tree = slide.shapes._spTree
    sp_tree.remove(rect._element)
    sp_tree.insert(2, rect._element)
    return rect


def add_footer(slide: Slide, page_num: int):
    """Small rounded-square 'V' watermark + page number, bottom-right."""
    size = Inches(0.35)
    x = SLIDE_W - Inches(0.9)
    y = SLIDE_H - Inches(0.7)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, size, size)
    box.adjustments[0] = 0.28
    box.fill.solid()
    box.fill.fore_color.rgb = BLUE
    box.line.fill.background()
    box.shadow.inherit = False
    tf = box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "V"
    set_run(run, font=HEAD_FONT, size=14, color=WHITE, bold=True)

    num_box = slide.shapes.add_textbox(x - Inches(0.5), y, Inches(0.45), size)
    num_tf = num_box.text_frame
    num_tf.word_wrap = False
    p2 = num_tf.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    run2 = p2.add_run()
    run2.text = str(page_num)
    set_run(run2, size=12, color=SECONDARY)


def add_title(slide: Slide, text: str, *, align=PP_ALIGN.LEFT, top=None, width=None):
    """Slide title in Google Sans Bold."""
    if top is None:
        top = Inches(0.6)
    if width is None:
        width = Inches(12.0)
    box = slide.shapes.add_textbox(Inches(0.7), top, width, Inches(1.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, font=HEAD_FONT, size=40, color=INK, bold=True)
    return box


def add_subhead(slide: Slide, text: str, *, top, align=PP_ALIGN.LEFT, size=20):
    """Secondary supporting line under a title."""
    box = slide.shapes.add_textbox(Inches(0.7), top, Inches(12.0), Inches(0.7))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, color=SECONDARY)
    return box


def add_rounded_container(
    slide: Slide, x, y, w, h, *, fill=BLUE, radius=0.18
) -> object:
    """A rounded-square container echoing addon/icon.png's geometry."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = radius
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_category_legend(slide: Slide, entries: list[tuple[str, str, str]], *, top):
    """A column of color-swatch + name + plain-language gloss rows.

    `entries` is (CATEGORY_COLORS key, display name, gloss). Reused across
    every deck that walks the 7 audit categories, so all of them read as one
    visual system.
    """
    for i, (key, name, gloss) in enumerate(entries):
        row_y = top + Inches(0.95) * i
        swatch = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1.0),
            row_y,
            Inches(0.5),
            Inches(0.5),
        )
        swatch.adjustments[0] = 0.3
        swatch.fill.solid()
        swatch.fill.fore_color.rgb = CATEGORY_COLORS[key]
        swatch.line.fill.background()
        swatch.shadow.inherit = False

        name_box = slide.shapes.add_textbox(
            Inches(1.8), row_y - Inches(0.05), Inches(4.3), Inches(0.6)
        )
        name_box.text_frame.word_wrap = True
        pn = name_box.text_frame.paragraphs[0]
        rn = pn.add_run()
        rn.text = name
        set_run(rn, font=HEAD_FONT, size=20, color=INK, bold=True)

        gloss_box = slide.shapes.add_textbox(
            Inches(6.2), row_y - Inches(0.02), Inches(6.3), Inches(0.6)
        )
        gloss_box.text_frame.word_wrap = True
        pg = gloss_box.text_frame.paragraphs[0]
        rg = pg.add_run()
        rg.text = gloss
        set_run(rg, size=18, color=SECONDARY, italic=True)


def add_notes(slide: Slide, lines: list[tuple[str | None, str, str | None]]):
    """Populate speaker notes verbatim from a PRESENTATION_PLAN's Section B.

    Each entry is (speaker, quoted_line, stage_direction) — stage_direction
    (a parenthetical) renders italic on its own trailing run, never spoken.
    """
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.clear()
    for i, (speaker, quote, stage) in enumerate(lines):
        p = notes_tf.paragraphs[0] if i == 0 else notes_tf.add_paragraph()
        if speaker:
            r = p.add_run()
            r.text = f"{speaker}: "
            set_run(r, size=14, color=INK, bold=True)
        r2 = p.add_run()
        r2.text = quote
        set_run(r2, size=14, color=INK)
        if stage:
            r3 = p.add_run()
            r3.text = f"  {stage}"
            set_run(r3, size=14, color=SECONDARY, italic=True)


def new_slide(prs: Presentation) -> Slide:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    return slide


# --------------------------------------------------------------------------
# Loop icon helpers (Observe / Decide / Act / Checkpoint)
# --------------------------------------------------------------------------


def add_icon_container(slide: Slide, cx, cy, size=None):
    if size is None:
        size = Inches(1.8)
    x = Emu(int(cx - size / 2))
    y = Emu(int(cy - size / 2))
    return add_rounded_container(slide, x, y, size, size, radius=0.22)


def _icon_scale(size) -> float:
    """Glyph offsets below are tuned for a 1.8in container; scale for others."""
    return size / Inches(1.8)


def add_observe_icon(slide: Slide, cx, cy, size=None):
    if size is None:
        size = Inches(1.8)
    add_icon_container(slide, cx, cy, size=size)
    k = _icon_scale(size)
    lens_d = int(Inches(0.7) * k)
    lens = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Emu(int(cx - Inches(0.5) * k)),
        Emu(int(cy - Inches(0.5) * k)),
        lens_d,
        lens_d,
    )
    lens.fill.background()
    lens.line.color.rgb = WHITE
    lens.line.width = Pt(5)
    lens.shadow.inherit = False

    handle = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Emu(int(cx + Inches(0.15) * k)),
        Emu(int(cy + Inches(0.15) * k)),
        int(Inches(0.55) * k),
        int(Inches(0.14) * k),
    )
    handle.rotation = 45
    handle.adjustments[0] = 0.5
    handle.fill.solid()
    handle.fill.fore_color.rgb = WHITE
    handle.line.fill.background()
    handle.shadow.inherit = False


def add_decide_icon(slide: Slide, cx, cy, size=None):
    if size is None:
        size = Inches(1.8)
    add_icon_container(slide, cx, cy, size=size)
    k = _icon_scale(size)
    fb = slide.shapes.build_freeform(
        Emu(int(cx - Inches(0.45) * k)), Emu(int(cy + Inches(0.05) * k))
    )
    fb.add_line_segments(
        [
            (Emu(int(cx - Inches(0.1) * k)), Emu(int(cy + Inches(0.4) * k))),
            (Emu(int(cx + Inches(0.5) * k)), Emu(int(cy - Inches(0.35) * k))),
        ],
        close=False,
    )
    shape = fb.convert_to_shape()
    shape.line.color.rgb = WHITE
    shape.line.width = Pt(max(3, 6 * k))
    shape.fill.background()
    shape.shadow.inherit = False


def add_act_icon(slide: Slide, cx, cy, size=None):
    if size is None:
        size = Inches(1.8)
    add_icon_container(slide, cx, cy, size=size)
    k = _icon_scale(size)
    body = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Emu(int(cx - Inches(0.55) * k)),
        Emu(int(cy - Inches(0.12) * k)),
        int(Inches(0.9) * k),
        int(Inches(0.22) * k),
    )
    body.rotation = -30
    body.adjustments[0] = 0.5
    body.fill.solid()
    body.fill.fore_color.rgb = WHITE
    body.line.fill.background()
    body.shadow.inherit = False

    tip = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        Emu(int(cx + Inches(0.3) * k)),
        Emu(int(cy - Inches(0.35) * k)),
        int(Inches(0.22) * k),
        int(Inches(0.22) * k),
    )
    tip.rotation = 60
    tip.fill.solid()
    tip.fill.fore_color.rgb = WHITE
    tip.line.fill.background()
    tip.shadow.inherit = False


def add_checkpoint_icon(slide: Slide, cx, cy, size=None):
    if size is None:
        size = Inches(1.8)
    add_icon_container(slide, cx, cy, size=size)
    k = _icon_scale(size)
    for dx in (-Inches(0.18) * k, Inches(0.05) * k):
        bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Emu(int(cx + dx)),
            Emu(int(cy - Inches(0.35) * k)),
            int(Inches(0.16) * k),
            int(Inches(0.7) * k),
        )
        bar.adjustments[0] = 0.5
        bar.fill.solid()
        bar.fill.fore_color.rgb = WHITE
        bar.line.fill.background()
        bar.shadow.inherit = False


# --------------------------------------------------------------------------
# Flow-diagram helpers (pipeline appendix slides)
# --------------------------------------------------------------------------


def add_flow_step(slide: Slide, top, text, sub, *, width=None, height=None):
    """One box in a pipeline diagram: bold title + a dim subline."""
    if width is None:
        width = Inches(8.0)
    if height is None:
        height = Inches(0.6)
    left = Emu(int((SLIDE_W - width) / 2))
    box = add_rounded_container(slide, left, top, width, height, radius=0.2)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = Pt(2)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = text
    set_run(r1, font=HEAD_FONT, size=16, color=WHITE, bold=True)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = sub
    set_run(r2, size=11, color=WHITE)
    return box


def add_flow_chevron(slide: Slide, center_y):
    """A small downward-pointing triangle marking flow direction."""
    w, h = Inches(0.26), Inches(0.1)
    tri = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        Emu(int((SLIDE_W - w) / 2)),
        Emu(int(center_y - h / 2)),
        w,
        h,
    )
    tri.rotation = 180
    tri.fill.solid()
    tri.fill.fore_color.rgb = BLUE
    tri.line.fill.background()
    tri.shadow.inherit = False
