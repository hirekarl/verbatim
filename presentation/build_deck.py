"""Generate the Saturday demo deck for Verbatim.

Standalone script, not part of the `verbatim` package: run with
`uv run --with python-pptx python presentation/build_deck.py`. Content is
transcribed verbatim from `presentation/PRESENTATION_PLAN.md` Sections A/B;
design language (colors, type, icon geometry) lives in `deck_common.py`,
shared with `build_deck_week2.py`, not repeated here.
"""

from __future__ import annotations

from pathlib import Path

from deck_common import (
    BLUE,
    HEAD_FONT,
    ICON_PATH,
    INK,
    SECONDARY,
    SLIDE_H,
    SLIDE_W,
    WHITE,
    add_act_icon,
    add_category_legend,
    add_checkpoint_icon,
    add_decide_icon,
    add_flow_chevron,
    add_flow_step,
    add_footer,
    add_notes,
    add_observe_icon,
    add_rounded_container,
    add_subhead,
    add_title,
    new_slide,
    set_run,
)
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.util import Emu, Inches

HERE = Path(__file__).parent
OUTPUT_PATH = HERE / "Verbatim - Saturday Demo.pptx"


# --------------------------------------------------------------------------
# Slide 1 — Title
# --------------------------------------------------------------------------


def build_slide_1():
    slide = new_slide(prs)
    if ICON_PATH.exists():
        slide.shapes.add_picture(
            str(ICON_PATH), Inches(5.867), Inches(1.6), Inches(1.6), Inches(1.6)
        )
    add_title(slide, "Verbatim", align=PP_ALIGN.CENTER, top=Inches(3.4))
    add_subhead(
        slide,
        "for copywriters",
        top=Inches(4.15),
        align=PP_ALIGN.CENTER,
        size=22,
    )
    add_subhead(
        slide,
        "20–30% of editing time → policing mechanical fixes",
        top=Inches(5.6),
        align=PP_ALIGN.CENTER,
        size=16,
    )
    add_footer(slide, 1)
    add_notes(
        slide,
        [
            ("KARL", '"I\'m Karl."', None),
            (
                "CHRISTINA",
                "\"I'm Christina. We built Verbatim for copywriters — the "
                "people who write a company's marketing emails, social posts, "
                'and web copy."',
                None,
            ),
            (
                "KARL",
                '"On a lot of marketing teams, the lead editor spends twenty '
                "to thirty percent of their editing time just policing "
                "mechanical stuff — voice, tone, formatting — instead of "
                "the actual writing and strategy. Right before a piece goes "
                "out, it needs that check. Verbatim automates it, so editors "
                'get their time back."',
                None,
            ),
        ],
    )


# --------------------------------------------------------------------------
# Slide 2 — Why two documents
# --------------------------------------------------------------------------


def add_document_glyph(slide: Slide, cx, cy):
    """Three thin white bars on a blue card, suggesting a document."""
    for i in range(3):
        bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Emu(int(cx - Inches(0.5))),
            Emu(int(cy + Inches(0.28) * i)),
            Inches(1.0),
            Inches(0.1),
        )
        bar.adjustments[0] = 0.5
        bar.fill.solid()
        bar.fill.fore_color.rgb = WHITE
        bar.line.fill.background()
        bar.shadow.inherit = False


def build_slide_2():
    slide = new_slide(prs)
    add_title(slide, "Why two documents")

    card_w, card_h = Inches(4.6), Inches(2.6)
    top = Inches(3.0)
    left1 = Inches(1.2)
    left2 = Inches(7.5)

    add_rounded_container(slide, left1, top, card_w, card_h)
    add_rounded_container(slide, left2, top, card_w, card_h)

    add_document_glyph(slide, int(left1 + card_w / 2), int(top + Inches(0.5)))
    add_document_glyph(slide, int(left2 + card_w / 2), int(top + Inches(0.5)))

    for card_left, label, gloss in (
        (left1, "Brand Style Guide", "permanent voice & rules"),
        (left2, "Campaign Brief", "specific to this one piece"),
    ):
        label_box = slide.shapes.add_textbox(
            card_left, top + Inches(1.5), card_w, Inches(0.5)
        )
        p = label_box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        set_run(r, font=HEAD_FONT, size=22, color=WHITE, bold=True)

        gloss_box = slide.shapes.add_textbox(
            card_left, top + Inches(2.0), card_w, Inches(0.5)
        )
        pg = gloss_box.text_frame.paragraphs[0]
        pg.alignment = PP_ALIGN.CENTER
        rg = pg.add_run()
        rg.text = gloss
        set_run(rg, size=14, color=WHITE)

    plus_box = slide.shapes.add_textbox(
        Inches(6.1), top + Inches(0.9), Inches(1.1), Inches(0.8)
    )
    pp = plus_box.text_frame.paragraphs[0]
    pp.alignment = PP_ALIGN.CENTER
    rp = pp.add_run()
    rp.text = "+"
    set_run(rp, font=HEAD_FONT, size=44, color=BLUE, bold=True)

    add_footer(slide, 2)
    add_notes(
        slide,
        [
            (
                "CHRISTINA",
                '"Every piece of copy gets checked against two different '
                "things. First, a brand style guide — a company's "
                "permanent voice and rules. Always sounds like this, never "
                "says that. The one we're demoing with today is built from "
                "Mailchimp's publicly published content style guide.\"",
                None,
            ),
            (
                "KARL",
                "\"Second, the campaign brief — what's specific to just "
                "this one piece? Who it's going to, what we want the reader "
                'to do, and by when."',
                None,
            ),
            (
                "CHRISTINA",
                "\"You need both. Copy that's perfectly on-brand but ignores "
                "what this campaign needs is useless. Copy that nails the "
                "ask but doesn't sound anything like the company is going "
                'to confuse people."',
                None,
            ),
            (
                "KARL",
                '"Checking both by hand means one person holding a style '
                "guide in one hand and a brief in the other, line by line. "
                "Two reviewers won't always catch the same things, and it's "
                "slow. That's the gap Verbatim closes.\"",
                None,
            ),
        ],
    )


# --------------------------------------------------------------------------
# Slide 3 — Loop: Observe
# --------------------------------------------------------------------------


def build_slide_3():
    slide = new_slide(prs)
    add_title(slide, "Loop — Observe", width=Inches(9.5))
    # Small badge beside the title, matching slide 4's fix — a full-size
    # icon stacked above the legend collided with the title text.
    add_observe_icon(slide, int(Inches(11.9)), int(Inches(1.2)), size=Inches(1.0))

    caption = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.9), Inches(11.0), Inches(0.5)
    )
    p = caption.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Docs/Drive API pull, then Christina's evaluator — regex, no model call"
    set_run(r, size=18, color=SECONDARY)

    add_category_legend(
        slide,
        [
            (
                "banned_words_and_competitors",
                "Banned words & competitors",
                "off-brand terms, no-go phrases",
            ),
            (
                "formatting_and_style",
                "Formatting & style",
                "commas, quotes, mechanical rules",
            ),
            (
                "channel_constraints",
                "Channel constraints",
                "character caps & rules per channel",
            ),
        ],
        top=Inches(2.6),
    )

    add_footer(slide, 3)
    add_notes(
        slide,
        [
            (
                "KARL",
                '"My half of the work involved getting the actual content '
                "onto the page. When a copywriter clicks ‘Run Verbatim "
                "Audit’ in the sidebar, with the brief and channel already "
                "filled in, Verbatim pulls the draft and the brief straight "
                "out of Google Docs — no copying and pasting into another "
                'tool."',
                None,
            ),
            (
                "CHRISTINA",
                "\"While that's happening, my evaluator runs over the text "
                "— a rules engine I built and tested against real sample "
                "copy. It checks what's black-and-white: banned words, "
                "formatting details like commas and quotation marks, and "
                "channel limits — things like Twitter's character cap, or "
                "an email subject line that's too long or too generic. None "
                "of that touches the model, it's just pattern matching.\"",
                None,
            ),
        ],
    )


# --------------------------------------------------------------------------
# Slide 4 — Loop: Decide (+ 4-category legend)
# --------------------------------------------------------------------------


def build_slide_4():
    slide = new_slide(prs)
    add_title(slide, "Loop — Decide", width=Inches(9.5))
    # Small badge beside the title, not stacked above the legend — a
    # full-size icon there collided with the title text (see
    # loop_decide_layout_problem.png).
    add_decide_icon(slide, int(Inches(11.9)), int(Inches(1.2)), size=Inches(1.0))

    add_category_legend(
        slide,
        [
            ("tone_drift", "Tone drift", "still sound like the brand?"),
            (
                "information_hierarchy",
                "Information hierarchy",
                "important stuff said first?",
            ),
            (
                "cta_cadence",
                "CTA cadence",
                "clear ask, shown at the right moment?",
            ),
            ("readability", "Readability", "actually easy to read?"),
        ],
        top=Inches(2.3),
    )

    add_footer(slide, 4)
    add_notes(
        slide,
        [
            (
                "CHRISTINA",
                '"Whatever my evaluator flags gets fed straight into the '
                "agent's system prompt — the instructions that shape "
                "everything it does — as evidence it can point to, instead "
                'of guessing at the same rules a second time."',
                None,
            ),
            (
                "KARL",
                '"That same system prompt is what makes the agent judge the '
                "four things that actually need judgment, not pattern "
                "matching. Does the writing still sound like the brand, or "
                "has the tone drifted? Is the important information said "
                "first, or is it buried three paragraphs down? Is there a "
                "clear CTA, and does it show up at the right moment and not "
                "too often? CTA means ‘call to action’ — the part of the "
                "copy that tells someone what to actually do next, like "
                "‘sign up’ or ‘buy now.’ And is it actually easy to read? "
                "Christina's findings, the brand guidelines, the document, "
                "and the brief all go into that one system prompt, and the "
                'agent decides what to flag."',
                None,
            ),
        ],
    )


# --------------------------------------------------------------------------
# Slide 5 — Loop: Act
# --------------------------------------------------------------------------


def build_slide_5():
    slide = new_slide(prs)
    add_title(slide, "Loop — Act")
    add_act_icon(slide, int(Inches(3.0)), int(Inches(3.9)))

    box = slide.shapes.add_textbox(Inches(5.2), Inches(3.3), Inches(7.0), Inches(1.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "create_suggestion (rewrite) or create_inline_comment (structural)"
    set_run(r, size=20, color=SECONDARY)

    add_footer(slide, 5)
    add_notes(
        slide,
        [
            (
                "KARL",
                '"From there the agent has two moves: propose a suggested '
                "edit for a rewrite, or leave a comment for anything more "
                "structural than a simple text swap. It keeps going until "
                "it's out of things to flag or hits a round limit, then "
                'prints a summary of what it posted."',
                None,
            ),
        ],
    )


# --------------------------------------------------------------------------
# Slide 6 — Checkpoint
# --------------------------------------------------------------------------


def build_slide_6():
    slide = new_slide(prs)
    add_title(slide, "Checkpoint")
    add_checkpoint_icon(slide, int(Inches(3.0)), int(Inches(3.9)))

    box = slide.shapes.add_textbox(Inches(5.2), Inches(3.5), Inches(7.0), Inches(0.8))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Reviewable. Never silent."
    set_run(r, font=HEAD_FONT, size=26, color=INK, bold=True)

    add_footer(slide, 6)
    add_notes(
        slide,
        [
            (
                "KARL",
                '"Nothing Verbatim does gets applied outright. Everything '
                "shows up as something reviewable — a suggested edit or a "
                'comment — never a silent change to the doc."',
                None,
            ),
            (
                "CHRISTINA",
                "\"That's because this copy is headed for a final sign-off "
                "— nobody wants their draft rewritten out from under them. "
                "The copywriter goes through every suggestion and comment "
                "one at a time and accepts or rejects each. Only what they "
                'accept moves forward."',
                None,
            ),
        ],
    )


# --------------------------------------------------------------------------
# Slide 7 — Live demo
# --------------------------------------------------------------------------


def build_slide_7():
    slide = new_slide(prs)
    add_title(slide, "Live demo", align=PP_ALIGN.CENTER, top=Inches(3.1))
    add_subhead(
        slide,
        "via the Verbatim Add-on sidebar",
        top=Inches(4.0),
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 7)
    add_notes(
        slide,
        [
            (
                "KARL",
                '"Let\'s see it in a doc."',
                "(open the Verbatim sidebar on the draft; the campaign "
                "brief link is already filled in)",
            ),
            (
                "CHRISTINA",
                '"We\'ll run it against Email."',
                '(pick Email from the channel dropdown, click "Run Verbatim Audit")',
            ),
            (
                "CHRISTINA",
                '"The evaluator and the model both run behind the scenes, '
                "so we'll click ‘Check Status’ until it's done thinking.\"",
                '(click "Check Status" once or twice until the results card appears)',
            ),
            (
                "KARL",
                '"Here\'s what came back."',
                "(read off the sidebar's category breakdown and counts — "
                "call out whatever actually shows up, not a fixed number)",
            ),
            (
                "CHRISTINA",
                '"And here in the doc—"',
                "(point at 2-3 of the actual suggestions or comments that "
                "landed, whatever they turn out to be)",
            ),
        ],
    )


# --------------------------------------------------------------------------
# Slide 8 — Close
# --------------------------------------------------------------------------


def build_slide_8():
    slide = new_slide(prs)
    if ICON_PATH.exists():
        slide.shapes.add_picture(
            str(ICON_PATH), Inches(5.867), Inches(2.2), Inches(1.6), Inches(1.6)
        )
    add_title(slide, "Verbatim", align=PP_ALIGN.CENTER, top=Inches(4.2))
    add_subhead(slide, "Thank you.", top=Inches(4.95), align=PP_ALIGN.CENTER)
    add_footer(slide, 8)
    add_notes(
        slide,
        [
            (
                "KARL",
                '"What you just saw is running for real — the same '
                "backend, the same Add-on, live in this doc. Not a "
                'mockup."',
                None,
            ),
            ("CHRISTINA", '"That\'s Verbatim. Thank you."', None),
        ],
    )


def build_slide_9():
    """Appendix: full pipeline, not part of the timed 5:00 script.

    Adapted from README.md's Mermaid flowchart, updated to the Add-on/Cloud
    Run path this deck actually demos (the README diagram documents the CLI
    entrypoint) and restyled in Verbatim's own palette/type rather than
    Mermaid's default theme, so it reads as one system with the rest of the
    deck.
    """
    slide = new_slide(prs)
    add_title(slide, "Under the hood")
    add_subhead(
        slide, "the full pipeline, for anyone curious", top=Inches(1.4), size=16
    )

    steps = [
        ("Run Verbatim Audit", "Add-on sidebar — brief & channel already set"),
        ("Fetch draft + brief", "Google Docs / Drive API"),
        ("Evaluator runs", "regex — banned words, formatting, channel limits"),
        ("LLM tool-calling loop", "create_suggestion / create_inline_comment"),
        ("Check Status", "Cloud Run backend — submit/poll, no 60s timeout"),
        ("Review in Google Docs", "accept / reject inline → sign-off"),
    ]
    box_h = Inches(0.55)
    gap = Inches(0.15)
    top0 = Inches(2.3)
    for i, (text, sub) in enumerate(steps):
        row_top = top0 + (box_h + gap) * i
        add_flow_step(slide, row_top, text, sub, height=box_h)
        if i > 0:
            prev_bottom = top0 + (box_h + gap) * (i - 1) + box_h
            add_flow_chevron(slide, Emu(int((prev_bottom + row_top) / 2)))

    add_footer(slide, 9)
    add_notes(
        slide,
        [
            (
                None,
                "Appendix — not part of the timed 5:00 script. Full "
                "technical pipeline for reference or Q&A; mirrors "
                "README.md's architecture diagram, updated to the live "
                "Add-on/Cloud Run path this deck actually demos.",
                None,
            ),
        ],
    )


prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

for builder in (
    build_slide_1,
    build_slide_2,
    build_slide_3,
    build_slide_4,
    build_slide_5,
    build_slide_6,
    build_slide_7,
    build_slide_8,
    build_slide_9,
):
    builder()

prs.save(str(OUTPUT_PATH))
print(f"Wrote {OUTPUT_PATH} ({len(prs.slides)} slides)")
