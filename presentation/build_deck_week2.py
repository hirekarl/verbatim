"""Generate the week 2 demo deck for Verbatim.

Standalone script, not part of the `verbatim` package: run with
`uv run --with python-pptx python presentation/build_deck_week2.py`. Content
is transcribed verbatim from `presentation/PRESENTATION_PLAN_WEEK2.md`
Sections A/B; design language (colors, type, icon geometry) lives in
`deck_common.py`, shared with `build_deck.py` so both decks read as one
visual system.
"""

from __future__ import annotations

from pathlib import Path

from deck_common import (
    BLUE,
    HEAD_FONT,
    ICON_PATH,
    INK,
    SECONDARY,
    SLIDE_H,
    SLIDE_W,
    WHITE,
    add_act_icon,
    add_decide_icon,
    add_footer,
    add_notes,
    add_observe_icon,
    add_rounded_container,
    add_subhead,
    add_title,
    new_slide,
    set_run,
)
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches

HERE = Path(__file__).parent
OUTPUT_PATH = HERE / "Verbatim - Week 2 Demo.pptx"

# --------------------------------------------------------------------------
# Slide 1 — Role & Pain Point
# --------------------------------------------------------------------------


def build_slide_1():
    slide = new_slide(prs)
    if ICON_PATH.exists():
        slide.shapes.add_picture(
            str(ICON_PATH), Inches(5.867), Inches(1.3), Inches(1.4), Inches(1.4)
        )
    add_title(slide, "Role & Pain Point", align=PP_ALIGN.CENTER, top=Inches(3.0))
    add_subhead(
        slide,
        "Verbatim, for copywriters",
        top=Inches(3.9),
        align=PP_ALIGN.CENTER,
        size=20,
    )
    add_subhead(
        slide,
        "20–30% of editing time → policing mechanical fixes",
        top=Inches(5.3),
        align=PP_ALIGN.CENTER,
        size=16,
    )
    add_footer(slide, 1)
    add_notes(
        slide,
        [
            (
                "CHRISTINA",
                "\"We're Karl and Christina, and Verbatim is still for "
                "copywriters — the people who write a company's marketing "
                'copy."',
                None,
            ),
            (
                "KARL",
                "\"The pain point hasn't changed either: twenty to thirty "
                "percent of editing time goes to policing mechanical stuff "
                "instead of real writing. This week was about making the "
                "agent that gets that time back actually hold up under real "
                'use."',
                None,
            ),
        ],
    )


# --------------------------------------------------------------------------
# Slide 2 — Week 5: The Original Agent
# --------------------------------------------------------------------------


def build_slide_2():
    slide = new_slide(prs)
    add_title(slide, "Week 5 — The Original Agent")

    row_y = Inches(3.3)
    spacing = Inches(3.6)
    left0 = Inches(3.0)
    for i, (icon_fn, label) in enumerate(
        [
            (add_observe_icon, "Observe"),
            (add_decide_icon, "Decide"),
            (add_act_icon, "Act"),
        ]
    ):
        cx = int(left0 + spacing * i)
        icon_fn(slide, cx, int(row_y), size=Inches(1.3))
        label_box = slide.shapes.add_textbox(
            Emu(cx - Inches(1.0)), row_y + Inches(1.0), Inches(2.0), Inches(0.5)
        )
        p = label_box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        set_run(r, font=HEAD_FONT, size=18, color=INK, bold=True)

    caption = slide.shapes.add_textbox(
        Inches(0.7), Inches(5.2), Inches(12.0), Inches(0.8)
    )
    p = caption.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = (
        "One agent, one loop, one checkpoint — judging all 4 subjective "
        "categories, backed by a deterministic evaluator for the rest."
    )
    set_run(r, size=18, color=SECONDARY)

    add_footer(slide, 2)
    add_notes(
        slide,
        [
            (
                "KARL",
                "\"Quick recap of where we started. Week 5's agent was one "
                "tool-calling loop, judging all four subjective categories "
                "at once, backed by a deterministic evaluator for the "
                "mechanical ones. It read the doc and the brief, decided "
                "what to flag, and posted suggestions or comments — nothing "
                'landed without an explicit accept in Google Docs."',
                None,
            ),
        ],
    )


# --------------------------------------------------------------------------
# Slide 3 — Week 6: Two Specialist Agents (Karl)
# --------------------------------------------------------------------------


def build_slide_3():
    slide = new_slide(prs)
    add_title(slide, "Week 6 — Two Specialist Agents")

    card_w, card_h = Inches(4.6), Inches(2.6)
    top = Inches(2.4)
    left1 = Inches(1.2)
    left2 = Inches(7.5)

    add_rounded_container(slide, left1, top, card_w, card_h)
    add_rounded_container(slide, left2, top, card_w, card_h)

    for card_left, label, gloss in (
        (left1, "Structural Agent", "info hierarchy, CTA cadence — comment only"),
        (left2, "Line-Editor Agent", "tone, readability — suggest or comment"),
    ):
        label_box = slide.shapes.add_textbox(
            card_left, top + Inches(0.8), card_w, Inches(0.6)
        )
        p = label_box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        set_run(r, font=HEAD_FONT, size=22, color=WHITE, bold=True)

        gloss_box = slide.shapes.add_textbox(
            card_left, top + Inches(1.4), card_w, Inches(0.8)
        )
        gloss_box.text_frame.word_wrap = True
        pg = gloss_box.text_frame.paragraphs[0]
        pg.alignment = PP_ALIGN.CENTER
        rg = pg.add_run()
        rg.text = gloss
        set_run(rg, size=14, color=WHITE)

    concurrent_box = slide.shapes.add_textbox(
        Inches(6.1), top + Inches(1.0), Inches(1.1), Inches(0.8)
    )
    pp = concurrent_box.text_frame.paragraphs[0]
    pp.alignment = PP_ALIGN.CENTER
    rp = pp.add_run()
    rp.text = "⇄"
    set_run(rp, font=HEAD_FONT, size=36, color=BLUE, bold=True)

    caption = slide.shapes.add_textbox(
        Inches(0.7), top + card_h + Inches(0.4), Inches(12.0), Inches(0.8)
    )
    p = caption.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = (
        "Concurrent threads, one tool each (a hard constraint) — an "
        "orchestrator merges results and flags cross-agent overlaps"
    )
    set_run(r, size=16, color=SECONDARY, italic=True)

    add_footer(slide, 3)
    add_notes(
        slide,
        [
            (
                "KARL",
                '"This week, one agent became two. The Structural agent '
                "judges information hierarchy and CTA cadence and can only "
                "leave comments; the Line-Editor agent judges tone and "
                "readability and can only propose rewrites — restricting "
                "each one to a single tool turned a soft preference into a "
                "hard constraint. They run concurrently on separate "
                "threads, and an orchestrator merges what they find — "
                "flagging, not resolving, any spot where both agents "
                "independently raise an issue. Narrower prompts fixed real "
                "false positives in each agent, and we made sure one "
                "agent's crash doesn't silently swallow the other's "
                'already-posted work."',
                None,
            ),
        ],
    )


# --------------------------------------------------------------------------
# Slide 4 — Week 6: Now Covering All 7 (Christina)
# --------------------------------------------------------------------------


def build_slide_4():
    slide = new_slide(prs)
    add_title(slide, "Week 6 — Now Covering All 7")

    card_w = Inches(5.4)
    card_h = Inches(4.3)
    top = Inches(2.1)
    left1 = Inches(0.9)
    left2 = Inches(7.0)

    add_rounded_container(slide, left1, top, card_w, card_h)
    add_rounded_container(slide, left2, top, card_w, card_h)

    structural_lines = ["Structural", "", "Information hierarchy", "CTA cadence"]
    line_editor_lines = [
        "Line-Editor",
        "",
        "Tone drift",
        "Readability",
        "— now also —",
        "Formatting & style",
        "Banned words & competitors",
        "Channel constraints",
    ]

    for card_left, lines in ((left1, structural_lines), (left2, line_editor_lines)):
        box = slide.shapes.add_textbox(
            card_left + Inches(0.4), top + Inches(0.3), card_w - Inches(0.8), card_h
        )
        tf = box.text_frame
        tf.word_wrap = True
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            r = p.add_run()
            r.text = line
            if i == 0:
                set_run(r, font=HEAD_FONT, size=22, color=WHITE, bold=True)
            elif line.startswith("—"):
                set_run(r, size=15, color=WHITE, italic=True)
            elif line:
                set_run(r, size=16, color=WHITE)

    caption = slide.shapes.add_textbox(
        Inches(0.7), top + card_h + Inches(0.2), Inches(12.0), Inches(0.8)
    )
    p = caption.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = (
        "PR #83 — Line-Editor now transcribes the evaluator's deterministic "
        "findings into real suggestions/comments, not just prompt context"
    )
    set_run(r, size=14, color=SECONDARY, italic=True)

    add_footer(slide, 4)
    add_notes(
        slide,
        [
            (
                "CHRISTINA",
                '"My half of this: I built the Structural agent from '
                "scratch, solo, against failing tests Karl wrote before he "
                "went out for a day — my first time in the agent code "
                "instead of the evaluator. By the back half of the week I "
                "was shipping fixes directly into files that used to be "
                "Karl-only — the orchestrator, the shared category "
                "validation, even the LLM client — because the split held "
                "up well enough to work across. The biggest one: the "
                "Line-Editor agent doesn't just read the evaluator's "
                "findings as context anymore, it transcribes them into "
                "real suggestions and comments. All seven categories reach "
                'the doc now, not four."',
                None,
            ),
        ],
    )


# --------------------------------------------------------------------------
# Slide 5 — Eval Evidence
# --------------------------------------------------------------------------


def build_slide_5():
    slide = new_slide(prs)
    add_title(slide, "Eval Evidence")

    entries = [
        (
            "Golden path — Mobile App Launch",
            "Caught a banned competitor comparison from two angles at once "
            "(structural hierarchy + line-editor banned-words, same "
            "overlapping span), plus real tone and readability fixes. "
            "7 findings, 5 categories, no errors.",
        ),
        (
            "Adversarial input — Feature Sunset",
            "A banned word repeated 5+ times, ageist language, two CTA "
            "stacks, heavy passive voice — hit the round cap and stopped "
            "early, but still returned 14 valid, correctly-categorized "
            "findings across all 7 categories instead of crashing.",
        ),
    ]
    top = Inches(2.3)
    for i, (label, body) in enumerate(entries):
        row_top = top + Inches(2.1) * i
        label_box = slide.shapes.add_textbox(
            Inches(0.9), row_top, Inches(11.5), Inches(0.5)
        )
        p = label_box.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = label
        set_run(r, font=HEAD_FONT, size=20, color=INK, bold=True)

        body_box = slide.shapes.add_textbox(
            Inches(0.9), row_top + Inches(0.55), Inches(11.5), Inches(1.4)
        )
        body_box.text_frame.word_wrap = True
        pb = body_box.text_frame.paragraphs[0]
        rb = pb.add_run()
        rb.text = body
        set_run(rb, size=16, color=SECONDARY)

    add_footer(slide, 5)
    add_notes(
        slide,
        [
            (
                "CHRISTINA",
                "\"We ran all six eval cards against this week's split. On "
                "the clean one — a mobile app launch email — it caught a "
                "banned competitor comparison from two angles at once: a "
                "structural hierarchy issue and a line-editor banned-words "
                "hit, flagged as the same overlapping span, plus real tone "
                'and readability fixes."',
                None,
            ),
            (
                "CHRISTINA",
                '"On the messiest one — a feature-sunset notice with a '
                "banned word repeated five times, ageist language, and two "
                "separate stacks of CTAs — it actually hit its round cap "
                "and had to stop early. But it still came back with "
                "fourteen valid, correctly-categorized findings instead of "
                'crashing or returning something broken."',
                None,
            ),
        ],
    )


# --------------------------------------------------------------------------
# Slide 6 — Live demo
# --------------------------------------------------------------------------


def build_slide_6():
    slide = new_slide(prs)
    add_title(slide, "Live Demo", align=PP_ALIGN.CENTER, top=Inches(3.1))
    add_subhead(
        slide,
        "via the Verbatim Add-on sidebar",
        top=Inches(4.0),
        align=PP_ALIGN.CENTER,
    )
    add_subhead(
        slide,
        "Eval Card 3 — Integration Announcement (Blog)",
        top=Inches(4.6),
        align=PP_ALIGN.CENTER,
        size=14,
    )
    add_footer(slide, 6)
    add_notes(
        slide,
        [
            (
                "KARL",
                "\"Let's run it — this is one of the six eval cards we just "
                'used to score the split."',
                "(open the Verbatim sidebar on the draft; the campaign "
                "brief link is already filled in)",
            ),
            (
                "CHRISTINA",
                '"Blog channel this time."',
                '(pick Blog, click "Run Verbatim Audit")',
            ),
            (
                "KARL",
                '"Two agents means two threads doing this at once now — '
                "we'll check status until both come back.\"",
                '(click "Check Status" once or twice until the results card appears)',
            ),
            (
                "CHRISTINA",
                '"Here\'s the breakdown."',
                "(read off the category counts — call out whatever actually shows up)",
            ),
            (
                "KARL",
                '"And here\'s what actually landed in the doc."',
                "(point at 2-3 of the real suggestions or comments that landed)",
            ),
        ],
    )


prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

for builder in (
    build_slide_1,
    build_slide_2,
    build_slide_3,
    build_slide_4,
    build_slide_5,
    build_slide_6,
):
    builder()

prs.save(str(OUTPUT_PATH))
print(f"Wrote {OUTPUT_PATH} ({len(prs.slides)} slides)")
