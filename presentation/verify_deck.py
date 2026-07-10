"""Structural sanity check for the generated Saturday demo deck.

Run with `uv run --with python-pptx python presentation/verify_deck.py`
after `build_deck.py`. Catches silent generation bugs (missing notes, wrong
slide count/titles) before ever opening PowerPoint. Not a substitute for the
manual visual/round-trip/timing checks in the plan.
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation

DECK_PATH = Path(__file__).parent / "Verbatim - Saturday Demo.pptx"

EXPECTED_TITLES = [
    "Verbatim",
    "Why two documents",
    "Loop — Observe",
    "Loop — Decide",
    "Loop — Act",
    "Checkpoint",
    "Live demo",
    "Verbatim",
    "Under the hood",
]

# Appendix slides aren't part of the timed 5:00 script, so their notes are a
# presenter reminder rather than a KARL:/CHRISTINA: speaker line.
APPENDIX_SLIDES = {9}


def slide_title(slide) -> str | None:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if text in EXPECTED_TITLES:
            return text
    return None


def main() -> int:
    if not DECK_PATH.exists():
        print(f"FAIL: {DECK_PATH} does not exist — run build_deck.py first")
        return 1

    prs = Presentation(str(DECK_PATH))
    errors = []

    if len(prs.slides) != len(EXPECTED_TITLES):
        errors.append(
            f"expected {len(EXPECTED_TITLES)} slides, found {len(prs.slides)}"
        )

    for i, slide in enumerate(prs.slides, start=1):
        expected_title = EXPECTED_TITLES[i - 1] if i <= len(EXPECTED_TITLES) else None
        title = slide_title(slide)
        if expected_title and title != expected_title:
            errors.append(
                f"slide {i}: expected title {expected_title!r}, got {title!r}"
            )

        notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
        if not notes.strip():
            errors.append(f"slide {i}: notes field is empty")
        elif (
            i not in APPENDIX_SLIDES
            and "KARL:" not in notes
            and "CHRISTINA:" not in notes
        ):
            errors.append(f"slide {i}: notes missing KARL:/CHRISTINA: speaker markers")

    if errors:
        print(f"FAIL: {len(errors)} issue(s) found in {DECK_PATH.name}")
        for e in errors:
            print(f"  - {e}")
        return 1

    print(f"OK: {DECK_PATH.name} — {len(prs.slides)} slides, all titled and noted")
    return 0


if __name__ == "__main__":
    sys.exit(main())
